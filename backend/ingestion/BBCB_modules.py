
import requests
import os
import gdown
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import zipfile
from PIL import Image
from moviepy import VideoFileClip
from transformers import (
    BlipProcessor, BlipForConditionalGeneration,
    Blip2Processor, Blip2ForConditionalGeneration,
    CLIPProcessor, CLIPModel,
    BertTokenizer, BertForSequenceClassification,
    BartTokenizer, BartForConditionalGeneration
)

MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY")
DEEPGRAM_API_KEY = os.getenv("DEEPGRAM_API_KEY")
LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")

class Parser:
    def __init__(self):
        self.raarr = open("order.txt","a+")
        self.AI_MODELS = self.initialize_models()

    def initialize_models(self):
    # 🤖 Initialize AI models for PDF processing
        self.raarr.write("0")
        """Initialize BLIP, BLIP2, CLIP, BERT, and BART models"""
        try:
            # BLIP for image captioning
            blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

            # BLIP2 for advanced image understanding
            '''try:
                blip2_processor = Blip2Processor.from_pretrained("Salesforce/blip2-opt-2.7b")
                blip2_model = Blip2ForConditionalGeneration.from_pretrained("Salesforce/blip2-opt-2.7b")
            except:
                blip2_processor, blip2_model = None, None
                print("⚠ BLIP2 not available, using BLIP only")'''

            # CLIP for image-text matching
            clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
            clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")

            # BERT for text classification/understanding
            bert_tokenizer = BertTokenizer.from_pretrained("bert-base-uncased")
            bert_model = BertForSequenceClassification.from_pretrained("bert-base-uncased")

            # BART for text summarization
            bart_tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn")
            bart_model = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn")

            return {
                'blip_processor': blip_processor,
                'blip_model': blip_model,
                #'blip2_processor': blip2_processor,
                #'blip2_model': blip2_model,
                'clip_processor': clip_processor,
                'clip_model': clip_model,
                'bert_tokenizer': bert_tokenizer,
                'bert_model': bert_model,
                'bart_tokenizer': bart_tokenizer,
                'bart_model': bart_model
            }
        except Exception as e:
            print(f"❌ Error initializing models: {e}")
            return None

    # 🤖 Enhanced AI question answering
    def ask_mistral_enhanced(self,context, question, file_type="unknown"):
        self.raarr.write("\n12")
        """Enhanced Mistral AI with context about file type"""
        try:
            url = "https://api.mistral.ai/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {MISTRAL_API_KEY}",
                "Content-Type": "application/json"
            }

            system_prompt = f"""You are a helpful assistant that answers questions based on the provided context from {file_type} files.
            The context may include text content, image descriptions, video transcripts, and AI-generated summaries.
            Provide comprehensive answers based on all available information."""

            payload = {
                "model": "mistral-large-latest",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Context: {context}\n\nQuestion: {question}"}
                ],
                "temperature": 0.7,
                "max_tokens": 1000
            }

            response = requests.post(url, headers=headers, json=payload, timeout=60)
            response.raise_for_status()

            return response.json()['choices'][0]['message']['content']
        except Exception as e:
            print(f"Error with Mistral API: {e}")
            return f"Error generating answer: {str(e)}"

    def match_segment(self,answer, segments):
        self.raarr.write("\n13")
        if not segments:
            return None

        answer_words = answer.lower().split()
        best_match = None
        best_score = 0

        for seg in segments:
            seg_text = seg["text"].lower()
            score = sum(1 for word in answer_words if word in seg_text)
            if score > best_score:
                best_score = score
                best_match = seg
        return best_match["start"] if best_match else None

    def process_document(self, uploaded_file, gdrive_url, question):
        self.raarr.write("\n16")
        """Enhanced document processing supporting multiple formats"""
        try:
            # Determine input source and file type
            file_path = None
            source_info = ""
            file_type = "unknown"

            if gdrive_url and gdrive_url.strip():
                # Process Google Drive URL
                print("🔗 Processing Google Drive URL...")
                downloaded_path, download_msg = self.download_from_google_drive(gdrive_url.strip())

                if downloaded_path:
                    file_path = downloaded_path
                    source_info = f"Downloaded from Google Drive: {download_msg}"
                    file_type = self.get_file_type_from_url(gdrive_url.strip())
                else:
                    return f"Failed to download from Google Drive: {download_msg}", [], "No video found", "Download failed", None, None

            elif uploaded_file:
                # Use uploaded file
                file_path = uploaded_file.name if hasattr(uploaded_file, 'name') else uploaded_file
                source_info = "File uploaded directly"

            else:
                return "No file uploaded or Google Drive URL provided", [], "No video found", "Please upload a file or provide a Google Drive URL", None, None

            if not question.strip():
                return f"File processed successfully ({source_info}), but no question provided", [], "No video found", "Please ask a question", None, None

            # Determine file type from extension
            file_ext = os.path.splitext(file_path)[1].lower()
            print(f"📄 Processing {file_ext} file: {file_path}")

            # Process based on file type
            videos = []
            images = []
            content_paths = []

            if file_ext == '.pdf':
                text, images = self.process_pdf(file_path)
                file_type = "PDF"
                # Get page files for matching
                content_paths = [f"components/pages/page_{i}.txt" for i in range(1, 100)
                            if os.path.exists(f"components/pages/page_{i}.txt")]

            elif file_ext == '.docx':
                text, images = self.process_docx(file_path)
                file_type = "Word Document"

            elif file_ext in ['.xlsx', '.xls']:
                text, images = self.process_excel(file_path)
                file_type = "Excel Spreadsheet"
                # Get sheet files for matching
                content_paths = [f"components/sheets/{f}" for f in os.listdir("components/sheets")
                            if f.endswith('.txt')]

            elif file_ext == '.pptx':
                text, content_paths = self.extract_text_and_slides(file_path)
                videos, slide_images = self.extract_media(file_path)
                images.extend(slide_images)
                file_type = "PowerPoint Presentation"

            else:
                return f"Unsupported file format: {file_ext}", [], "Unsupported format", "Please upload a supported file format", None, None

            # Process video if available (mainly for PPTX)
            transcript = "No video found"
            segments = []
            if videos:
                video = videos[0]
                transcript, segments = self.transcribe_with_deepgram(video)
                context = f"{file_type} Content:\n{text}\n\nVideo Transcript:\n{transcript}"
            else:
                context = f"{file_type} Content:\n{text}"

            # Get answer from Mistral AI
            answer = self.ask_mistral_enhanced(context, question, file_type)

            # Add source information to answer
            full_answer = f"[Source: {source_info}]\n[File Type: {file_type}]\n\n{answer}"

            # Find matching timestamp and create video clip
            video_clip = None
            if videos and segments:
                timestamp = self.match_segment(answer, segments)
                if timestamp is not None:
                    video_clip = self.extract_video_clip(videos[0], timestamp)

            # Find matching content
            matched_content = self.find_matching_content(answer, content_paths)

            return text, images, transcript, full_answer, video_clip, matched_content

        except Exception as e:
            error_msg = f"Error processing request: {str(e)}"
            print(error_msg)
            return error_msg, [], "Error", "Error", None, None



    def extract_video_clip(self, video_path, center_time, output_path="components/clips/clip.mp4"):
        self.raarr.write("\n14")
        try:
            clip = VideoFileClip(video_path)
            start = max(center_time - 30, 0)
            end = min(center_time + 30, clip.duration)

            if end - start < 10:
                start = max(0, end - 60)

            subclip = clip.subclip(start, end)
            subclip.write_videofile(
                output_path,
                codec="libx264",
                audio_codec="aac",
                verbose=False,
                logger=None,
                temp_audiofile="temp-audio.m4a",
                remove_temp=True
            )

            clip.close()
            subclip.close()

            return output_path
        except Exception as e:
            print(f"Error extracting video clip: {e}")
            return None


    def transcribe_with_deepgram(self, video_path):
        self.raarr.write("\n11")
        try:
            headers = {
                "Authorization": f"Token {DEEPGRAM_API_KEY}",
                "Content-Type": "video/mp4"
            }

            with open(video_path, "rb") as f:
                response = requests.post(
                    "https://api.deepgram.com/v1/listen?model=nova-2&smart_format=true",
                    headers=headers,
                    data=f,
                    timeout=300
                )

            if response.status_code != 200:
                print(f"Deepgram API error: {response.status_code} - {response.text}")
                return "Transcription failed", []

            result = response.json()
            transcript = ""
            segments = []

            channels = result.get("results", {}).get("channels", [])
            if channels:
                alternatives = channels[0].get("alternatives", [])
                if alternatives:
                    words = alternatives[0].get("words", [])
                    for word in words:
                        segments.append({
                            "start": word.get("start", 0),
                            "end": word.get("end", 0),
                            "text": word.get("word", "")
                        })
                        transcript += word.get("word", "") + " "

            return transcript.strip(), segments
        except Exception as e:
            print(f"Error transcribing video: {e}")
            return "Transcription failed", []


    def extract_media(self, file_path):
        self.raarr.write("\n10")
        """Extract media from PPTX files"""
        videos = []
        images = []
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                for file in z.namelist():
                    if file.startswith("ppt/media/"):
                        ext = os.path.splitext(file)[-1].lower()
                        filename = os.path.basename(file)
                        data = z.read(file)

                        if ext in ['.mp4', '.avi', '.mov', '.wmv']:
                            path = f"components/videos/{filename}"
                            with open(path, "wb") as f:
                                f.write(data)
                            videos.append(path)
                        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                            path = f"components/images/{filename}"
                            with open(path, "wb") as f:
                                f.write(data)
                            images.append(path)
        except Exception as e:
            print(f"Error extracting media: {e}")

        return videos, images


    def extract_text_and_slides(self, pptx_file_path):
        self.raarr.write("9")
        try:
            prs = Presentation(pptx_file_path)
            text_data = []
            slide_paths = []

            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                text_data.append(slide_text.strip())

                # Save slide as a text file
                slide_file_path = f"components/slides/slide_{i+1}.txt"
                with open(slide_file_path, "w", encoding="utf-8") as f:
                    f.write(slide_text)
                slide_paths.append(slide_file_path)

            final_text = "\n\n".join(text_data)
            with open("components/texts/all_text.txt", "w", encoding="utf-8") as f:
                f.write(final_text)
            return final_text, slide_paths
        except Exception as e:
            print(f"Error extracting text and slides: {e}")
            return "Error extracting slides", []


    def process_excel(self, excel_path):
        self.raarr.write("8")
        """Process Excel file"""
        workbook = openpyxl.load_workbook(excel_path)
        text_data = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = f"Sheet: {sheet_name}\n"

            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text += row_text + "\n"

            text_data.append(sheet_text)

            # Save individual sheet
            sheet_file = f"components/sheets/sheet_{sheet_name}.txt"
            with open(sheet_file, "w", encoding="utf-8") as f:
                f.write(sheet_text)

        full_text = "\n\n".join(text_data)

        # Save full text
        with open("components/texts/all_text.txt", "w", encoding="utf-8") as f:
            f.write(full_text)

        return full_text, []



    # 📝 Word document processing
    def process_docx(self, docx_path):
        self.raarr.write("7")
        """Process Word document"""
        try:
            doc = docx.Document(docx_path)
            text_data = []
            images = []

            # Extract text
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_data.append(paragraph.text)

            # Extract images from document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob
                    img_path = f"components/images/docx_image_{len(images) + 1}.png"
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    images.append(img_path)

            full_text = "\n".join(text_data)

            # Save text
            with open("components/texts/all_text.txt", "w", encoding="utf-8") as f:
                f.write(full_text)

            return full_text, images

        except Exception as e:
            print(f"Error processing DOCX: {e}")
            return f"Error processing DOCX: {str(e)}", []



    def process_pdf(self, pdf_path):
        self.raarr.write("6")
        """Process PDF using PyMuPDF with BLIP, BLIP2, CLIP, BERT, and BART"""
        try:
            doc = fitz.open(pdf_path)
            text_data = []
            images = []
            image_descriptions = []

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)

                # Extract text
                page_text = page.get_text()
                if page_text.strip():
                    text_data.append(f"Page {page_num + 1}:\n{page_text}")

                    # Save page text
                    page_file = f"components/pages/page_{page_num + 1}.txt"
                    with open(page_file, "w", encoding="utf-8") as f:
                        f.write(f"Page {page_num + 1}:\n{page_text}")

                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)

                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_path = f"components/images/pdf_page_{page_num + 1}img{img_index + 1}.png"
                            pix.save(img_path)
                            images.append(img_path)

                            # Analyze image with AI models
                            img_description = self.analyze_image_with_ai(img_path)
                            image_descriptions.append(f"Image from page {page_num + 1}: {img_description}")

                        pix = None
                    except Exception as e:
                        print(f"Error extracting image {img_index} from page {page_num + 1}: {e}")

            doc.close()

            # Combine all text
            full_text = "\n\n".join(text_data)

            # Add image descriptions to text
            if image_descriptions:
                full_text += "\n\n--- Image Descriptions ---\n" + "\n".join(image_descriptions)

            # Use BART for summarization if text is very long
            if len(full_text) > 5000 and self.AI_MODELS:
                try:
                    bart_inputs = self.AI_MODELS['bart_tokenizer'](full_text[:1024], return_tensors="pt",
                                                            max_length=1024, truncation=True)
                    summary_ids = self.AI_MODELS['bart_model'].generate(bart_inputs["input_ids"],
                                                                max_length=200, min_length=50,
                                                                length_penalty=2.0, num_beams=4)
                    bart_summary = self.AI_MODELS['bart_tokenizer'].decode(summary_ids[0], skip_special_tokens=True)
                    full_text = f"BART Summary:\n{bart_summary}\n\n--- Full Content ---\n{full_text}"
                except Exception as e:
                    print(f"BART summarization failed: {e}")

            # Save full text
            with open("components/texts/all_text.txt", "w", encoding="utf-8") as f:
                f.write(full_text)

            return full_text, images

        except Exception as e:
            print(f"Error processing PDF: {e}")
            return f"Error processing PDF: {str(e)}", []



    def download_from_google_drive(self, url, output_dir="downloads"):
        self.raarr.write("4")
        """Enhanced download function supporting multiple formats"""
        try:
            file_id = self.extract_file_id_from_url(url)
            if not file_id:
                return None, "Could not extract file ID from URL"

            file_type = self.get_file_type_from_url(url)

            if '/folders/' in url:
                # Handle folder download
                folder_url = f"https://drive.google.com/drive/folders/{file_id}"
                print(f"📁 Downloading folder: {folder_url}")

                try:
                    gdown.download_folder(folder_url, output=output_dir, quiet=False, use_cookies=False)

                    # Find supported files in downloaded folder
                    supported_files = []
                    for root, dirs, files in os.walk(output_dir):
                        for file in files:
                            ext = os.path.splitext(file)[1].lower()
                            if ext in ['.pptx', '.pdf', '.docx', '.xlsx', '.xls']:
                                supported_files.append(os.path.join(root, file))

                    if supported_files:
                        return supported_files[0], f"Downloaded folder and found {len(supported_files)} supported file(s)"
                    else:
                        return None, "No supported files found in the downloaded folder"

                except Exception as e:
                    return None, f"Error downloading folder: {str(e)}"

            else:
                # Handle individual file download
                if file_type == 'slides':
                    export_url = f"https://docs.google.com/presentation/d/{file_id}/export/pptx"
                    output_file = os.path.join(output_dir, f"google_slides_{file_id}.pptx")
                elif file_type == 'docs':
                    export_url = f"https://docs.google.com/document/d/{file_id}/export/docx"
                    output_file = os.path.join(output_dir, f"google_docs_{file_id}.docx")
                elif file_type == 'sheets':
                    export_url = f"https://docs.google.com/spreadsheets/d/{file_id}/export/xlsx"
                    output_file = os.path.join(output_dir, f"google_sheets_{file_id}.xlsx")
                else:
                    # Try as regular file
                    file_url = f"https://drive.google.com/uc?id={file_id}"
                    output_file = os.path.join(output_dir, f"downloaded_file_{file_id}")

                    try:
                        gdown.download(file_url, output_file, quiet=False)
                        return output_file, "File downloaded successfully"
                    except Exception as e:
                        return None, f"Error downloading file: {str(e)}"

                print(f"📄 Downloading {file_type}: {export_url}")

                try:
                    response = requests.get(export_url, timeout=60)
                    response.raise_for_status()

                    with open(output_file, "wb") as f:
                        f.write(response.content)

                    if os.path.exists(output_file) and os.path.getsize(output_file) > 0:
                        return output_file, f"Google {file_type} downloaded successfully"
                    else:
                        return None, f"Downloaded {file_type} file is empty"

                except requests.exceptions.RequestException as e:
                    return None, f"Error downloading Google {file_type}: {str(e)}. Make sure the file is publicly accessible."

        except Exception as e:
            return None, f"Error processing Google Drive URL: {str(e)}"




    # 🔍 Enhanced matching functions
    def match_segment(self, answer, segments):
        self.raarr.write("\n13")
        if not segments:
            return None

        answer_words = answer.lower().split()
        best_match = None
        best_score = 0

        for seg in segments:
            seg_text = seg["text"].lower()
            score = sum(1 for word in answer_words if word in seg_text)
            if score > best_score:
                best_score = score
                best_match = seg

        return best_match["start"] if best_match else None

    def analyze_image_with_ai(self,image_path, question=None):
        self.raarr.write("5")
        """Analyze image using BLIP, BLIP2, and CLIP models"""
        if not self.AI_MODELS:
            return "AI models not available"

        try:
            image = Image.open(image_path).convert('RGB')

            # Use BLIP for image captioning
            blip_inputs = self.AI_MODELS['blip_processor'](image, return_tensors="pt")
            blip_out = self.AI_MODELS['blip_model'].generate(**blip_inputs, max_length=50)
            blip_caption = self.AI_MODELS['blip_processor'].decode(blip_out[0], skip_special_tokens=True)

            # Use BLIP2 if available
            blip2_caption = ""
            if self.AI_MODELS['blip2_processor'] and self.AI_MODELS['blip2_model']:
                try:
                    blip2_inputs = self.AI_MODELS['blip2_processor'](image, return_tensors="pt")
                    blip2_out = self.AI_MODELS['blip2_model'].generate(**blip2_inputs, max_length=50)
                    blip2_caption = self.AI_MODELS['blip2_processor'].decode(blip2_out[0], skip_special_tokens=True)
                except:
                    blip2_caption = "BLIP2 analysis failed"

            # Combine captions
            full_caption = f"BLIP: {blip_caption}"
            if blip2_caption:
                full_caption += f"\nBLIP2: {blip2_caption}"

            return full_caption

        except Exception as e:
            return f"Error analyzing image: {str(e)}"

    def find_matching_content(self,answer, content_paths):
        self.raarr.write("\n15")
        if not content_paths:
            return None

        best_content = None
        best_score = 0
        answer_words = answer.lower().split()

        for path in content_paths:
            try:
                with open(path, "r", encoding="utf-8") as f:
                    content = f.read().lower()
                score = sum(1 for word in answer_words if word in content)
                if score > best_score:
                    best_score = score
                    best_content = path
            except Exception as e:
                print(f"Error reading content file {path}: {e}")

        return best_content


if __name__ == "__main__":
    parser = Parser()
    # Example usage
    text, images, transcript, answer, video_clip, matched_content = parser.process_document(
        uploaded_file = "final_code/vectorDB/M1 - Why Startups Fail.pdf",
        gdrive_url = None,
        question = "What is the main topic of this document?"
    )
    print(answer)
    if video_clip:
        print(f"Video clip created at: {video_clip}")
    if matched_content:
        print(f"Matched content: {matched_content}")
    if text:
        print(f"Extracted text: {text[:500]}...")
    if images:  
        print(f"Extracted {len(images)} images.")
