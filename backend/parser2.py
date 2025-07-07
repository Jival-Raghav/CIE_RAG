
import os
import uuid
import json
import fitz  # PyMuPDF
import torch
from PIL import Image
from typing import List, Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
from sentence_transformers import SentenceTransformer
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import Document, VectorStoreIndex  # adapt if needed

from mistralai.client import MistralClient
#from mistralai.models.chat_completion import ChatMessage

import hashlib

# Constants
STORAGE_DIR = "image_storage"
os.makedirs(STORAGE_DIR, exist_ok=True)

class SentenceSplitter:
    """Splits large text into overlapping chunks."""
    def __init__(self, chunk_size=250, chunk_overlap=50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def split_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap."""
        words = text.split()
        chunks = []
        start = 0
        while start < len(words):
            end = start + self.chunk_size
            chunk = words[start:end]
            chunks.append(" ".join(chunk))
            start += self.chunk_size - self.chunk_overlap
        return chunks

class Parser:
    """Chatbot that uses Qdrant for storage and retrieval of text/image knowledge."""
    def __init__(self, mistral_api_key: str, qdrant_url: str = "http://localhost:6333"):
        self.mistral_client = MistralClient(
            api_key=mistral_api_key,
            # endpoint="https://api.mistral.ai/v1"  # changed hear 
        )
        self.qdrant_url = qdrant_url
        self.embedding_model = HuggingFaceEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")
        self.embedding_dim = 384
        self.qdrant_client = QdrantClient(url=qdrant_url)
        self.collection_name = "Trial_For_parsersa"

        # Load BLIP model for image captioning
        from transformers import BlipProcessor, BlipForConditionalGeneration
        self.blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        self.blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
        
        # try blip -2 
        # from transformers import Blip2Processor, Blip2ForConditionalGeneration
        # self.blip_processor = Blip2Processor.from_pretrained("Salesforce/blip2-opt-2.7b-coco")
        # self.blip_model = Blip2ForConditionalGeneration.from_pretrained("Salesforce/blip2-opt-2.7b-coco")

        # Create Qdrant collection if not exists
        collections = self.qdrant_client.get_collections().collections
        if self.collection_name not in [c.name for c in collections]:
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=VectorParams(size=self.embedding_dim, distance=Distance.COSINE)
            )

    def _store_image_and_metadata(self, img_data: bytes, img_ext: str, metadata: dict) -> str:
        """Saves image and its metadata to disk. only for debuggin and checking capiton"""
        img_hash = hashlib.sha256(img_data).hexdigest()[:16]
        img_id = f"{img_hash}_{uuid.uuid4().hex[:8]}"
        img_path = os.path.join(STORAGE_DIR, f"{img_id}.{img_ext}")
        with open(img_path, "wb") as f:
            f.write(img_data)
        with open(os.path.join(STORAGE_DIR, f"{img_id}.json"), "w") as f:
            json.dump(metadata, f)
        return img_id

    def generate_image_caption(self, image_path: str) -> str:
        """Generates caption for an image using BLIP -2 image captioning model."""
        if not hasattr(self, 'blip_model') or not hasattr(self, 'blip_processor'):
            return "Image captioning model not loaded."

        try:
            image = Image.open(image_path).convert("RGB")
            inputs = self.blip_processor(image, return_tensors="pt")
            with torch.no_grad():
                output = self.blip_model.generate(**inputs, max_length=50)
                # blip 2
                # caption = self.blip_processor.decode(output[0], skip_special_tokens=True)
                # for blip 
                caption = self.blip_processor.decode(output[0], skip_special_tokens=True)
            return caption
        except Exception as e:
            return f"Error generating caption: {e}" 
    

    def extract_text_from_pptx(self, pptx_path: str) -> List[Dict[str, Any]]:
        """Extracts text and images from PPTX, returns list of chunks and captions."""
        prs = Presentation(pptx_path)
        documents = []
        splitter = SentenceSplitter()

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                full_text = "\n".join(slide_text)
                chunks = splitter.split_text(full_text)
                for idx, chunk in enumerate(chunks):
                    documents.append({
                        "id": f"slide_{slide_num}_chunk_{idx}",
                        "text": chunk,
                        "metadata": {
                            "slide": slide_num,
                            "chunk_id": idx,
                            "type": "text"
                        }
                    })

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    img_bytes = image.blob
                    img_ext = image.ext
                    img_id = self._store_image_and_metadata(img_bytes, img_ext, {"slide": slide_num})
                    img_path = os.path.join(STORAGE_DIR, f"{img_id}.{img_ext}")
                    caption = self.generate_image_caption(img_path)
                    json_path = os.path.join(STORAGE_DIR, f"{img_id}.json")
                    if os.path.exists(json_path):
                        with open(json_path, "r") as jf:
                            meta = json.load(jf)
                        meta["caption"] = caption
                        with open(json_path, "w") as jf:
                            json.dump(meta, jf)
                    documents.append({
                        "id": f"slide_{slide_num}_image_{img_id}",
                        "text": caption,
                        "metadata": {
                            "slide": slide_num,
                            "image_id": img_id,
                            "type": "image",
                            "caption": caption,
                            "ext": img_ext
                        }
                    })

        return documents

    """ changed this to add video chunking """
    def extract_text_from_video(self, video_path: str) -> list:
        """
        Transcribes video to text and splits into chunks.
        Returns a list of dicts with chunked text and metadata.
        """
        splitter = SentenceSplitter()
        # Use your video class (import or pass as needed)
        from video_demo_1 import video  # adjust import as needed
        #full_text= video.youtube_video_to_text("video_path", STORAGE_DIR)
        full_text = video.local_video_to_text(video_path, STORAGE_DIR)
        if not full_text.strip():
            return []
        chunks = splitter.split_text(full_text)
        documents = []
        for idx, chunk in enumerate(chunks):
            documents.append({
                "id": f"video_chunk_{idx}",
                "text": chunk,
                "metadata": {
                    "chunk_id": idx,
                    "type": "video_text",
                    "file_name": os.path.basename(video_path),
                    "source": "local"  # as now file is local you may need to change it as you passs from other function
                }
            })
        return documents
    
    
    
    def extract_text_from_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """Extracts text and images with captions from PDF file."""
        doc = fitz.open(pdf_path)
        splitter = SentenceSplitter()
        documents = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                chunks = splitter.split_text(text)
                for idx, chunk in enumerate(chunks):
                    documents.append({
                        "id": f"page_{page_num}_chunk_{idx}",
                        "text": chunk,
                        "metadata": {
                            "page": page_num,
                            "chunk_id": idx,
                            "type": "text"
                        }
                    })

            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                img_id = self._store_image_and_metadata(img_bytes, img_ext, {"page": page_num})
                img_path = os.path.join(STORAGE_DIR, f"{img_id}.{img_ext}")
                caption = self.generate_image_caption(img_path)
                json_path = os.path.join(STORAGE_DIR, f"{img_id}.json")
                if os.path.exists(json_path):
                    with open(json_path, "r") as jf:
                        meta = json.load(jf)
                    meta["caption"] = caption
                    with open(json_path, "w") as jf:
                        json.dump(meta, jf)
                documents.append({
                    "id": f"page_{page_num}_image_{img_id}",
                    "text": caption,
                    "metadata": {
                        "page": page_num,
                        "image_id": img_id,
                        "type": "image",
                        "caption": caption,
                        "ext": img_ext
                    }
                })

        doc.close()
        return documents

    def add_knowledge(self, documents: List[Dict[str, Any]]):
        """Embeds and adds documents (text + captions) to Qdrant."""
        texts = [doc.text for doc in documents]
        # Use HuggingFaceEmbedding's get_text_embedding method
        embeddings = [self.embedding_model.get_text_embedding(text) for text in texts]
        points = []
        for doc, emb in zip(documents, embeddings):
            points.append(PointStruct(
                id=str(uuid.uuid4()),
                vector=emb,
                payload={**doc.metadata, "text": doc.text}
            ))
        self.qdrant_client.upsert(collection_name=self.collection_name, points=points)

    def extract_and_index_files(self, directory_path: str):
        """
        Extracts content from all files in a directory, creates Document objects, manages image metadata, and builds the index.
        """
        all_documents = []
        image_metadata = {}

        # Get all files from the directory
        if not os.path.exists(directory_path):
            print(f"Directory not found: {directory_path}")
            return None

        file_paths = []
        for file_name in os.listdir(directory_path):
            file_path = os.path.join(directory_path, file_name)
            if os.path.isfile(file_path) and file_name.lower().endswith(('.pdf', '.pptx', '.mp4', '.avi', '.mov', '.mkv')):
                file_paths.append(file_path)

        if not file_paths:
            print(f"No supported files found in directory: {directory_path}")
            return None

        for file_path in file_paths:
            file_name = os.path.basename(file_path)
            print(f"Processing {file_name}...")

            try:
                if file_name.lower().endswith('.pptx'):
                    content = self.extract_text_from_pptx(file_path)
                elif file_name.lower().endswith('.pdf'):
                    content = self.extract_text_from_pdf(file_path)
                elif file_name.lower().endswith(('.mp4', '.avi', '.mov', '.mkv')):
                    content = self.extract_text_from_video(file_path)
                else:
                    continue

                for item in content:
                    meta = item.get("metadata", {})
                    doc_type = meta.get("type", "text")

                    if doc_type in ("text", "video_text"):
                        doc = Document(
                            text=item["text"],
                            metadata={
                                **meta,
                                "file_name": file_name,
                                "source": "local"
                            }
                        )
                        all_documents.append(doc)
                    elif doc_type == "image":
                        doc = Document(
                            text=f"Image caption: {item['text']}",
                            metadata={
                                **meta,
                                "file_name": file_name,
                                "source": "local"
                            }
                        )
                        all_documents.append(doc)

                        image_metadata[meta["image_id"]] = {
                            "caption": meta["caption"],
                            "file_name": file_name,
                            "image_path": os.path.join(STORAGE_DIR, f"{meta['image_id']}.{meta.get('ext', 'png')}"),
                            "ext": meta.get("ext", "png"),
                            "source": "local",
                            **({"slide": meta["slide"]} if "slide" in meta else {}),
                            **({"page": meta["page"]} if "page" in meta else {}),
                        }

            except Exception as e:
                print(f"Error processing {file_name}: {str(e)}")
                continue

        print("Indexing documents...")

        if not all_documents:
            print("No content extracted from the files.")
            return None
        ####################################
        #if "image_metadata" not in st.session_state:
         #   st.session_state.image_metadata = {}
        ####################################
        # Add to Qdrant
        self.add_knowledge(all_documents)
        
        # (Optional) Build and return a local index if you still need it for other purposes
        # embed_model = HuggingFaceEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")
        # index = VectorStoreIndex.from_documents(all_documents, embed_model=embed_model)
        print("✅ Processing complete!")
        print(f"Indexed {len(all_documents)} chunks from {len(file_paths)} files")

        # return index


def main():
    """Main entry point for command-line use of the parser."""
    mistral_key = ""  # For security, use env vars in production
    qdrant_url = "http://localhost:6333"

    parser = Parser(mistral_api_key=mistral_key, qdrant_url=qdrant_url)
    
    # Example local file paths - modify these to your actual file paths
    """ file_paths =  "directory"
    """
    file_paths = "parsing\data"
    
    parser.extract_and_index_files(file_paths)
    print("Local files processed and indexed.")


if __name__ == "__main__":
    main()

"""# Example usage in your script (not CLI):
file_paths = ["path/to/your/document.pdf", "path/to/your/presentation.pptx"]
parser = Parser(mistral_api_key="Put_yout_mistral_key")
parser.extract_and_index_files(file_paths)
"""



"""
QDRANT CLOUD MIGRATION GUIDE
============================

To switch from local Qdrant to Qdrant Cloud, make these changes:

1. REQUIRED CHANGES IN CODE:
   --------------------------
   (line is not known exactly )
   A) Line 32 - Change default parameter:
      FROM: qdrant_url: str = "http://localhost:6333"
      TO:   qdrant_url: str = "https://your-cluster-url.qdrant.tech"
   
   B) Line 38 - Add API key to QdrantClient:
      FROM: self.qdrant_client = QdrantClient(url=qdrant_url)
      TO:   self.qdrant_client = QdrantClient(
                url=qdrant_url,
                api_key="your_qdrant_api_key_here"
            )
   
   C) Line 171 - Update main() function:
      FROM: qdrant_url = "http://localhost:6333"
      TO:   qdrant_url = "https://your-cluster-url.qdrant.tech"
   
   D) Line 173 - Add API key parameter (modify __init__ to accept it):
      Add new parameter: qdrant_api_key: str = None
      
      Then modify QdrantClient initialization to:
      self.qdrant_client = QdrantClient(
          url=qdrant_url,
          api_key=qdrant_api_key
      )

2. REQUIRED CREDENTIALS:
   ---------------------
   - Qdrant Cloud URL (from your cluster dashboard)
   - Qdrant API Key (from your cluster settings)

3. RECOMMENDED CHANGES:
   --------------------
   - Store credentials in environment variables:
     QDRANT_URL = "https://your-cluster-url.qdrant.tech"
     QDRANT_API_KEY = "your_api_key_here"
   
   - Use os.getenv() to load them:
     qdrant_url = os.getenv("QDRANT_URL", "http://localhost:6333")
     qdrant_api_key = os.getenv("QDRANT_API_KEY")

4. UPDATED CONSTRUCTOR SIGNATURE:
   -------------------------------
   def __init__(self, mistral_api_key: str, qdrant_url: str = "https://your-cluster.qdrant.tech", qdrant_api_key: str = None):

5. UPDATED MAIN FUNCTION:
   ----------------------
   qdrant_url = "https://your-cluster-url.qdrant.tech"
   qdrant_api_key = "your_qdrant_api_key_here"
   parser = Parser(mistral_api_key=mistral_key, qdrant_url=qdrant_url, qdrant_api_key=qdrant_api_key)

SUMMARY: 4 lines to change + 1 new parameter to add for cloud migration.
"""
